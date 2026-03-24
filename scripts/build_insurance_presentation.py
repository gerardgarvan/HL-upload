"""Build UF-branded PowerPoint presentation from insurance summary tables (Phase 7+8).

Reads Phase 5, Phase 6, and Phase 8 CSV outputs, produces:
- reports/insurance_tables_YYYY-MM-DD.pptx (14-slide PowerPoint presentation)

Presentation structure:
- Slide 1: Title slide with cohort sizes
- Slide 2: Overview table (Phase 5 - all patients)
- Slide 3: Combined post-treatment (Phase 6 - all patients)
- Slide 4: Chemotherapy insurance (Phase 5)
- Slide 5: Chemotherapy post-treatment (Phase 6)
- Slide 6: Radiation insurance (Phase 5)
- Slide 7: Radiation post-treatment (Phase 6)
- Slide 8: SCT insurance (Phase 5)
- Slide 9: SCT post-treatment (Phase 6)
- Slide 10: DX enrollment comparison (Phase 8)
- Slide 11: Chemotherapy enrollment comparison (Phase 8)
- Slide 12: Radiation enrollment comparison (Phase 8)
- Slide 13: SCT enrollment comparison (Phase 8)
- Slide 14: Unknown post-treatment encounter breakdown (Phase 8)

All tables use UF Health branding (blue #003087, orange #FA4616) with native
PowerPoint tables (not embedded images) for editability.

Usage:
    python scripts/build_insurance_presentation.py [config/paths.toml]
"""

import sys
from pathlib import Path
from datetime import datetime

PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))

import polars as pl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from src.load.config import load_and_validate_config

# UF Health brand colors
UF_BLUE = RGBColor.from_string("003087")
UF_ORANGE = RGBColor.from_string("FA4616")
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(0xCC, 0xD5, 0xEA)  # Light blue tint for odd rows
LIGHT_ORANGE = RGBColor(0xFD, 0xD9, 0xCC)  # Light orange tint for even rows
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)  # Near-black for body text

# Standard payer category order (9 categories)
PAYER_CATEGORY_ORDER = [
    "Medicare",
    "Medicaid",
    "Dual eligible",
    "Private",
    "Other government",
    "Self-pay",
    "Other",
    "Unavailable",
    "Unknown"
]


def _extract_cohort_size(df: pl.DataFrame, column_name: str) -> int:
    """Extract cohort size by summing the N column.

    Args:
        df: DataFrame with payer category rows and N/% columns
        column_name: Base column name (e.g., "Primary Insurance", "Post-Treatment Insurance")

    Returns:
        Total N across all payer categories
    """
    n_col = f"{column_name} (N)"
    if n_col not in df.columns:
        return 0
    return df[n_col].sum()


def _read_csv_safe(path: Path) -> pl.DataFrame | None:
    """Read CSV file with error handling.

    Args:
        path: Path to CSV file

    Returns:
        DataFrame if successful, None if file missing or read error
    """
    try:
        if not path.exists():
            print(f"  [WARNING] CSV not found: {path}")
            return None
        return pl.read_csv(path)
    except Exception as e:
        print(f"  [WARNING] Error reading {path}: {e}")
        return None


def _add_table_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    df: pl.DataFrame,
    n_pct_columns: list[str],
    payer_col: str = "Payer Category"
) -> None:
    """Add a table slide to the presentation.

    Args:
        prs: Presentation object
        title: Slide title
        subtitle: Slide subtitle (includes cohort size N=X)
        df: DataFrame with payer categories and N_Pct columns
        n_pct_columns: List of column names ending in (N_Pct) to display
        payer_col: Name of payer category column
    """
    # Use blank layout (index 6) for full control
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title text box
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(22)
    title_para.font.bold = True
    title_para.font.color.rgb = UF_BLUE
    title_para.alignment = PP_ALIGN.LEFT

    # Subtitle text box
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.65), Inches(9.0), Inches(0.3)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(12)
    subtitle_para.font.italic = True
    subtitle_para.font.color.rgb = DARK_TEXT
    subtitle_para.alignment = PP_ALIGN.LEFT

    # Table positioning
    rows = len(df) + 1  # Data rows + header
    cols = 1 + len(n_pct_columns)  # Payer Category + data columns

    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.1), Inches(9.0), Inches(4.0)
    ).table

    # Header row
    header_cells = [table.cell(0, i) for i in range(cols)]

    # First header cell (Payer Category)
    cell = header_cells[0]
    cell.fill.solid()
    cell.fill.fore_color.rgb = UF_BLUE
    cell.text = "Payer Category"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = cell.text_frame.paragraphs[0]
    para.font.size = Pt(11)
    para.font.bold = True
    para.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER

    # Data column headers (remove (N_Pct) suffix)
    for i, col_name in enumerate(n_pct_columns, start=1):
        cell = header_cells[i]
        cell.fill.solid()
        cell.fill.fore_color.rgb = UF_BLUE
        # Remove " (N_Pct)" suffix for cleaner header
        display_name = col_name.replace(" (N_Pct)", "")
        cell.text = display_name
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(df.iter_rows(named=True)):
        table_row_idx = row_idx + 1  # Skip header

        # Alternate row colors (light blue for odd, light orange for even)
        row_color = LIGHT_BLUE if row_idx % 2 == 0 else LIGHT_ORANGE

        # Payer Category cell
        cell = table.cell(table_row_idx, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_color
        cell.text = str(row[payer_col])
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.bold = True
        para.font.color.rgb = DARK_TEXT
        para.alignment = PP_ALIGN.LEFT

        # Data cells
        for col_idx, col_name in enumerate(n_pct_columns, start=1):
            cell = table.cell(table_row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            cell.text = str(row[col_name]) if row[col_name] is not None else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = DARK_TEXT
            para.alignment = PP_ALIGN.CENTER

    # Set column widths
    table.columns[0].width = Inches(2.2)  # Payer Category column
    remaining_width = 9.0 - 2.2
    data_col_width = remaining_width / len(n_pct_columns)
    for i in range(1, cols):
        table.columns[i].width = Inches(data_col_width)


def _create_title_slide(
    prs: Presentation,
    total_n: int,
    chemo_n: int,
    radiation_n: int,
    sct_n: int
) -> None:
    """Create title slide with cohort sizes.

    Args:
        prs: Presentation object
        total_n: Total cohort size
        chemo_n: Chemotherapy cohort size
        radiation_n: Radiation cohort size
        sct_n: SCT cohort size
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.0)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Insurance Coverage by Treatment Type"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = UF_BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Colored accent bar
    accent = slide.shapes.add_shape(
        1,  # Rectangle shape type
        Inches(0.5), Inches(2.0), Inches(4.0), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = UF_ORANGE
    accent.line.fill.background()

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2), Inches(9.0), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Hodgkin Lymphoma Cohort — UF Health"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_TEXT
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Cohort sizes
    cohort_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.0), Inches(9.0), Inches(1.5)
    )
    cohort_frame = cohort_box.text_frame
    cohort_frame.text = f"Total Cohort: N = {total_n:,}"

    # Add remaining lines as new paragraphs
    for line in [
        f"Chemotherapy: N = {chemo_n:,}",
        f"Radiation: N = {radiation_n:,}",
        f"Stem Cell Transplant: N = {sct_n:,}"
    ]:
        p = cohort_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

    # Format first paragraph
    cohort_frame.paragraphs[0].font.size = Pt(14)
    cohort_frame.paragraphs[0].font.color.rgb = DARK_TEXT
    cohort_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def main(config_path: str | None = None) -> None:
    """Generate PowerPoint presentation from Phase 5 and 6 CSV outputs.

    Args:
        config_path: Optional path to config file (uses default if None)
    """
    print("=" * 60)
    print("INSURANCE TABLES POWERPOINT GENERATION")
    print("=" * 60)

    # Load config
    paths = load_and_validate_config(config_path)

    # Define input directories
    phase5_dir = PROJECT_ROOT / "reports" / "insurance_by_treatment"
    phase6_dir = PROJECT_ROOT / "reports" / "post_treatment_insurance"

    # Phase 5 CSV files (4 tables)
    overview_path = phase5_dir / "overview_table.csv"
    chemo_path = phase5_dir / "chemotherapy_table.csv"
    radiation_path = phase5_dir / "radiation_table.csv"
    sct_path = phase5_dir / "sct_table.csv"

    # Phase 6 CSV files (4 tables)
    combined_post_path = phase6_dir / "combined_post_treatment.csv"
    chemo_post_path = phase6_dir / "chemo_post_treatment.csv"
    radiation_post_path = phase6_dir / "radiation_post_treatment.csv"
    sct_post_path = phase6_dir / "sct_post_treatment.csv"

    # Phase 8 CSV files (5 tables)
    phase8_dir = PROJECT_ROOT / "reports" / "insurance_enr_comparison"
    dx_enr_path = phase8_dir / "dx_enr_comparison.csv"
    chemo_enr_path = phase8_dir / "chemo_enr_comparison.csv"
    radiation_enr_path = phase8_dir / "radiation_enr_comparison.csv"
    sct_enr_path = phase8_dir / "sct_enr_comparison.csv"
    unknown_breakdown_path = phase8_dir / "unknown_post_tx_encounter_breakdown.csv"

    print(f"\nReading CSV files...")
    print(f"  Phase 5 dir: {phase5_dir}")
    print(f"  Phase 6 dir: {phase6_dir}")
    print(f"  Phase 8 dir: {phase8_dir}")

    # Read Phase 5 CSVs
    overview_df = _read_csv_safe(overview_path)
    chemo_df = _read_csv_safe(chemo_path)
    radiation_df = _read_csv_safe(radiation_path)
    sct_df = _read_csv_safe(sct_path)

    # Read Phase 6 CSVs
    combined_post_df = _read_csv_safe(combined_post_path)
    chemo_post_df = _read_csv_safe(chemo_post_path)
    radiation_post_df = _read_csv_safe(radiation_post_path)
    sct_post_df = _read_csv_safe(sct_post_path)

    # Read Phase 8 CSVs
    dx_enr_df = _read_csv_safe(dx_enr_path)
    chemo_enr_df = _read_csv_safe(chemo_enr_path)
    radiation_enr_df = _read_csv_safe(radiation_enr_path)
    sct_enr_df = _read_csv_safe(sct_enr_path)
    unknown_breakdown_df = _read_csv_safe(unknown_breakdown_path)

    # Check if we have any data
    all_dfs = [
        overview_df, chemo_df, radiation_df, sct_df,
        combined_post_df, chemo_post_df, radiation_post_df, sct_post_df,
        dx_enr_df, chemo_enr_df, radiation_enr_df, sct_enr_df, unknown_breakdown_df
    ]
    available_count = sum(1 for df in all_dfs if df is not None)

    if available_count == 0:
        print("\n  [FATAL] No CSV files found. Cannot generate presentation.")
        print("          Run Phase 5 and 6 scripts first:")
        print("          - python scripts/build_insurance_by_treatment.py")
        print("          - python scripts/build_post_treatment_insurance.py")
        sys.exit(1)

    print(f"\n  Found {available_count} of 8 CSV files")

    # Extract cohort sizes for title slide
    total_n = 0
    chemo_n = 0
    radiation_n = 0
    sct_n = 0

    if overview_df is not None:
        # Detect primary insurance column (could be "Primary Insurance" or "First Diagnosis")
        n_pct_cols = [c for c in overview_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            first_col = n_pct_cols[0]
            base_name = first_col.replace(" (N_Pct)", "")
            total_n = _extract_cohort_size(overview_df, base_name)

    if chemo_df is not None:
        chemo_n = _extract_cohort_size(chemo_df, "Primary Insurance")

    if radiation_df is not None:
        radiation_n = _extract_cohort_size(radiation_df, "Primary Insurance")

    if sct_df is not None:
        sct_n = _extract_cohort_size(sct_df, "Primary Insurance")

    print(f"\nCohort sizes:")
    print(f"  Total: {total_n:,}")
    print(f"  Chemotherapy: {chemo_n:,}")
    print(f"  Radiation: {radiation_n:,}")
    print(f"  SCT: {sct_n:,}")

    # Create presentation
    print(f"\nBuilding PowerPoint presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    slide_count = 0

    # Slide 1: Title slide
    _create_title_slide(prs, total_n, chemo_n, radiation_n, sct_n)
    slide_count += 1
    print(f"  Slide {slide_count}: Title slide")

    # Slide 2: Overview table (Phase 5)
    if overview_df is not None:
        n_pct_cols = [c for c in overview_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Insurance Coverage Overview",
                f"All enrolled patients — N = {total_n:,}",
                overview_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: Overview table")

    # Slide 3: Combined Post-Treatment (Phase 6)
    if combined_post_df is not None:
        combined_n = _extract_cohort_size(combined_post_df, "Post-Treatment Insurance")
        _add_table_slide(
            prs,
            "Post-Treatment Insurance — All Patients",
            f"Most prevalent payer after last treatment date — N = {combined_n:,}",
            combined_post_df,
            ["Post-Treatment Insurance (N_Pct)"]
        )
        slide_count += 1
        print(f"  Slide {slide_count}: Combined post-treatment")

    # Slide 4: Chemotherapy Insurance (Phase 5)
    if chemo_df is not None:
        n_pct_cols = [c for c in chemo_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Chemotherapy Insurance",
                f"Insurance at primary, first, and last chemotherapy — N = {chemo_n:,}",
                chemo_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: Chemotherapy insurance")

    # Slide 5: Chemotherapy Post-Treatment (Phase 6)
    if chemo_post_df is not None:
        chemo_post_n = _extract_cohort_size(chemo_post_df, "Post-Treatment Insurance")
        _add_table_slide(
            prs,
            "Chemotherapy Post-Treatment Insurance",
            f"Most prevalent payer after last treatment date — N = {chemo_post_n:,}",
            chemo_post_df,
            ["Post-Treatment Insurance (N_Pct)"]
        )
        slide_count += 1
        print(f"  Slide {slide_count}: Chemotherapy post-treatment")

    # Slide 6: Radiation Insurance (Phase 5)
    if radiation_df is not None:
        n_pct_cols = [c for c in radiation_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Radiation Insurance",
                f"Insurance at primary, first, and last radiation — N = {radiation_n:,}",
                radiation_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: Radiation insurance")

    # Slide 7: Radiation Post-Treatment (Phase 6)
    if radiation_post_df is not None:
        radiation_post_n = _extract_cohort_size(radiation_post_df, "Post-Treatment Insurance")
        _add_table_slide(
            prs,
            "Radiation Post-Treatment Insurance",
            f"Most prevalent payer after last treatment date — N = {radiation_post_n:,}",
            radiation_post_df,
            ["Post-Treatment Insurance (N_Pct)"]
        )
        slide_count += 1
        print(f"  Slide {slide_count}: Radiation post-treatment")

    # Slide 8: SCT Insurance (Phase 5)
    if sct_df is not None:
        n_pct_cols = [c for c in sct_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Stem Cell Transplant Insurance",
                f"Insurance at primary, first, and last SCT — N = {sct_n:,}",
                sct_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: SCT insurance")

    # Slide 9: SCT Post-Treatment (Phase 6)
    if sct_post_df is not None:
        sct_post_n = _extract_cohort_size(sct_post_df, "Post-Treatment Insurance")
        _add_table_slide(
            prs,
            "SCT Post-Treatment Insurance",
            f"Most prevalent payer after last treatment date — N = {sct_post_n:,}",
            sct_post_df,
            ["Post-Treatment Insurance (N_Pct)"]
        )
        slide_count += 1
        print(f"  Slide {slide_count}: SCT post-treatment")

    # Slide 10: DX Enrollment Comparison (Phase 8)
    if dx_enr_df is not None:
        n_pct_cols = [c for c in dx_enr_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Diagnosis — Insurance by Enrollment Coverage",
                "Payer at first HL diagnosis: patients with vs without enrollment covering \u00b130 day window",
                dx_enr_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: DX enrollment comparison")

    # Slide 11: Chemo Enrollment Comparison (Phase 8)
    if chemo_enr_df is not None:
        n_pct_cols = [c for c in chemo_enr_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Chemotherapy — Insurance by Enrollment Coverage",
                "Payer at first/last chemo: patients with vs without enrollment covering \u00b130 day window",
                chemo_enr_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: Chemo enrollment comparison")

    # Slide 12: Radiation Enrollment Comparison (Phase 8)
    if radiation_enr_df is not None:
        n_pct_cols = [c for c in radiation_enr_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Radiation — Insurance by Enrollment Coverage",
                "Payer at first/last radiation: patients with vs without enrollment covering \u00b130 day window",
                radiation_enr_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: Radiation enrollment comparison")

    # Slide 13: SCT Enrollment Comparison (Phase 8)
    if sct_enr_df is not None:
        n_pct_cols = [c for c in sct_enr_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "SCT — Insurance by Enrollment Coverage",
                "Payer at first/last SCT: patients with vs without enrollment covering \u00b130 day window",
                sct_enr_df,
                n_pct_cols
            )
            slide_count += 1
            print(f"  Slide {slide_count}: SCT enrollment comparison")

    # Slide 14: Unknown Post-Treatment Encounter Breakdown (Phase 8)
    if unknown_breakdown_df is not None:
        n_pct_cols = [c for c in unknown_breakdown_df.columns if c.endswith(" (N_Pct)")]
        if n_pct_cols:
            _add_table_slide(
                prs,
                "Unknown Post-Treatment Payer — Encounter Breakdown",
                "Patients with Unknown post-treatment payer: how many encounters exist after last treatment?",
                unknown_breakdown_df,
                n_pct_cols,
                payer_col="Encounter Count Bin"
            )
            slide_count += 1
            print(f"  Slide {slide_count}: Unknown post-treatment breakdown")

    # Save presentation
    output_dir = PROJECT_ROOT / "reports"
    output_dir.mkdir(exist_ok=True)

    date_str = datetime.now().strftime("%Y-%m-%d")
    output_path = output_dir / f"insurance_tables_{date_str}.pptx"

    prs.save(str(output_path))

    print(f"\n{'=' * 60}")
    print(f"PRESENTATION COMPLETE")
    print(f"{'=' * 60}")
    print(f"  Slides created: {slide_count}")
    print(f"  Output: {output_path}")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    config_path = sys.argv[1] if len(sys.argv) > 1 else None
    main(config_path)
